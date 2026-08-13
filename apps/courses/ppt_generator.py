"""
PPT 自动生成工具 - 使用 python-pptx 从课程资料生成 PPT
"""
import os
import io
from django.conf import settings

# python-pptx 是可选依赖，运行时检测
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def generate_course_ppt(course, materials, output_path=None):
    """根据课程信息和资料生成 PPT

    返回 (file_path, file_name) 或 (None, error_message)
    """
    if not HAS_PPTX:
        return None, 'python-pptx 未安装，请执行: pip install python-pptx'

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== 第 1 页：封面 ==========
    slide = _add_blank_slide(prs)
    _add_colored_bg(slide, RGBColor(0x1A, 0x56, 0xDB))

    # 标题
    _add_text_box(slide, Inches(1), Inches(2), Inches(11.3), Inches(1.5),
                  course.title, font_size=Pt(40), color=RGBColor(0xFF, 0xFF, 0xFF),
                  bold=True, alignment=PP_ALIGN.CENTER)

    # 副标题
    subtitle = f'培训课程 · {course.category.name if course.category else "通用"}'
    _add_text_box(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.6),
                  subtitle, font_size=Pt(20), color=RGBColor(0xCC, 0xDD, 0xFF),
                  alignment=PP_ALIGN.CENTER)

    # 创建者 & 日期
    _add_text_box(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
                  f'创建者：{course.creator}    |    {course.created_at.strftime("%Y年%m月%d日")}',
                  font_size=Pt(14), color=RGBColor(0xAA, 0xBB, 0xDD),
                  alignment=PP_ALIGN.CENTER)

    # ========== 第 2 页：课程概述 ==========
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, '课程概述', RGBColor(0x1A, 0x56, 0xDB))

    _add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
                  course.description or '（无课程描述）',
                  font_size=Pt(18), color=RGBColor(0x33, 0x33, 0x33))

    # ========== 第 3 页：课程大纲 ==========
    slide = _add_blank_slide(prs)
    _add_slide_title(slide, '课程大纲', RGBColor(0x1A, 0x56, 0xDB))

    content_lines = []
    for i, m in enumerate(materials, 1):
        icon = {'ppt': 'PPT', 'word': '文档', 'pdf': 'PDF', 'video': '视频'}.get(m.file_type, '资料')
        content_lines.append(f'{i}. [{icon}] {m.title}')
        if m.description:
            content_lines.append(f'     {m.description[:80]}')

    _add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
                  '\n'.join(content_lines) if content_lines else '暂无课程资料',
                  font_size=Pt(18), color=RGBColor(0x33, 0x33, 0x33))

    # ========== 第 4 页之后：每个资料各一页摘要 ==========
    for i, m in enumerate(materials):
        slide = _add_blank_slide(prs)
        _add_slide_title(slide, m.title, RGBColor(0x1A, 0x56, 0xDB))

        # 文件类型标签
        type_colors = {
            'ppt': RGBColor(0xD0, 0x44, 0x23),
            'word': RGBColor(0x2B, 0x57, 0x9A),
            'pdf': RGBColor(0xE7, 0x4C, 0x3C),
            'video': RGBColor(0x28, 0xA7, 0x45),
            'image': RGBColor(0x6F, 0x42, 0xC1),
        }
        tag_color = type_colors.get(m.file_type, RGBColor(0x6C, 0x75, 0x7D))

        _add_shape_box(slide, Inches(0.8), Inches(1.5), Inches(2.5), Inches(0.5),
                       m.get_file_type_display(), fill_color=tag_color,
                       font_color=RGBColor(0xFF, 0xFF, 0xFF), font_size=Pt(14))

        desc = m.description or '（无描述信息）'
        _add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11.5), Inches(4.0),
                      desc, font_size=Pt(16), color=RGBColor(0x44, 0x44, 0x44))

        # 页脚
        _add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
                      f'资料 {i+1}/{len(materials)}    |    文件大小：{m.get_file_size_display()}',
                      font_size=Pt(11), color=RGBColor(0x99, 0x99, 0x99))

    # ========== 最后一页：结束 ==========
    slide = _add_blank_slide(prs)
    _add_colored_bg(slide, RGBColor(0x1A, 0x56, 0xDB))
    _add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                  '感谢学习', font_size=Pt(44), color=RGBColor(0xFF, 0xFF, 0xFF),
                  bold=True, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.8),
                  f'共 {len(materials)} 份学习资料，请按顺序完成学习',
                  font_size=Pt(20), color=RGBColor(0xCC, 0xDD, 0xFF),
                  alignment=PP_ALIGN.CENTER)

    # 保存
    if output_path is None:
        output_path = os.path.join(settings.MEDIA_ROOT, 'generated_ppts', str(course.pk))
    os.makedirs(output_path, exist_ok=True)

    file_name = f'{course.title}_课件.pptx'
    file_path = os.path.join(output_path, file_name)
    prs.save(file_path)

    return file_path, file_name


# ========== 辅助函数 ==========

def _add_blank_slide(prs):
    """添加空白幻灯片"""
    blank_layout = prs.slide_layouts[6]  # 空白布局
    return prs.slides.add_slide(blank_layout)


def _add_colored_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_slide_title(slide, title_text, accent_color):
    """添加页面标题栏"""
    # 左侧色条
    _add_shape_box(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5),
                   '', fill_color=accent_color)

    # 标题文字
    _add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                  title_text, font_size=Pt(32), color=accent_color, bold=True)


def _add_text_box(slide, left, top, width, height, text,
                  font_size=Pt(18), color=RGBColor(0x33, 0x33, 0x33),
                  bold=False, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment

    return txBox


def _add_shape_box(slide, left, top, width, height, text,
                   fill_color=None, font_color=RGBColor(0xFF, 0xFF, 0xFF),
                   font_size=Pt(14)):
    """添加矩形色块"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.CENTER

    return shape